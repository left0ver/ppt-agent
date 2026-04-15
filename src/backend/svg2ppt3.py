import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import cairosvg

def batch_svg_to_pptx(svg_directory, output_pptx):
    """
    将指定目录下的所有 SVG 文件批量保存到一个 PPTX 文件中。
    每张幻灯片对应一个 SVG 图像。
    """
    # 1. 初始化 PPT 演示文稿
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6] # 索引 6 为空白布局
    
    # 2. 获取目录下所有的 SVG 文件
    svg_dir_path = Path(svg_directory)
    svg_files = list(svg_dir_path.glob("*.svg"))
    
    if not svg_files:
        print(f"在目录 '{svg_directory}' 中没有找到 SVG 文件。")
        return

    print(f"找到 {len(svg_files)} 个 SVG 文件，开始处理...")
    
    # 定义一个临时的 PNG 文件名，用于中转
    temp_png_path = "temp_batch_convert.png"
    
    # 3. 遍历所有 SVG 文件
    for index, svg_file in enumerate(svg_files, start=1):
        try:
            print(f"[{index}/{len(svg_files)}] 正在处理: {svg_file.name}")
            
            # 将当前 SVG 转为高分辨率 PNG
            cairosvg.svg2png(url=str(svg_file), write_to=temp_png_path)
            
            # 添加一张新幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 插入生成的 PNG 图片
            # 这里设置左侧和顶部留白 1 英寸，宽度设为 8 英寸（根据标准 16:9 幻灯片尺寸可自行调整）
            slide.shapes.add_picture(temp_png_path, Inches(1), Inches(1), width=Inches(8))
            
        except Exception as e:
            print(f"处理文件 {svg_file.name} 时出错: {e}")
            continue

    # 4. 保存最终的 PPTX 文件
    prs.save(output_pptx)
    print(f"\n全部处理完成！文件已保存至: {output_pptx}")
    
    # 5. 清理临时文件
    if os.path.exists(temp_png_path):
        os.remove(temp_png_path)

# ==========================================
# 运行示例
# ==========================================
if __name__ == "__main__":
    # 请将这里的路径替换为你存放 SVG 文件的实际文件夹路径
    # 例如: "./my_svg_folder" 或 "C:/Users/name/Desktop/svgs"
    TARGET_DIRECTORY = "backend/final_ppt" 
    OUTPUT_FILE = "Batch_Output.pptx"
    
    batch_svg_to_pptx(TARGET_DIRECTORY, OUTPUT_FILE)
