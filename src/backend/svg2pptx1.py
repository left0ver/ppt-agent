from pptx import Presentation
from svg2pptx import SVGConverter, svg_to_pptx, Config
import os
from pathlib import Path
# Convert SVG file to PowerPoint
# svg_to_pptx("icon.svg", "output.pptx")


# Add SVG shapes to an existing presentation
prs = Presentation()

config = Config(
    curve_tolerance=0.5,
)
svg_dir = "backend/final_ppt"
converter = SVGConverter(config=config)
for file in os.listdir("backend/final_ppt"):
    # print(file)
    svg_path = Path(svg_dir) / file
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    converter.add_to_slide(svg_path, slide)
    print(svg_path)

prs.save("combined.pptx")
