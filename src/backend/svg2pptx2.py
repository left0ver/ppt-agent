from __future__ import annotations

import argparse
import re
import subprocess
import sys
import warnings
from pathlib import Path
from tempfile import TemporaryDirectory
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

PAGE_PATTERN = re.compile(r"page_(\d+)\.svg$", re.IGNORECASE)
PX_PER_INCH = 96
EMU_PER_INCH = 914400


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert SVG pages into a PPTX by exporting SVG to EMF with Inkscape."
    )
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=Path("final_ppt"),
        help="Directory containing page_*.svg files. Defaults to ./final_ppt.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("final_ppt.pptx"),
        help="Output PPTX path. Defaults to ./final_ppt.pptx.",
    )
    parser.add_argument(
        "--keep-emf-dir",
        type=Path,
        default=None,
        help="Optional directory to keep exported EMF files. If omitted, a temporary directory is used.",
    )
    return parser.parse_args()


def svg_sort_key(path: Path) -> tuple[int, str]:
    match = PAGE_PATTERN.search(path.name)
    if match:
        return int(match.group(1)), path.name
    return sys.maxsize, path.name


def list_svg_files(input_dir: Path) -> list[Path]:
    svg_files = [path for path in input_dir.glob("*.svg") if path.is_file()]
    svg_files.sort(key=svg_sort_key)
    return svg_files


def parse_numeric(value: str) -> float:
    match = re.match(r"^\s*([0-9]+(?:\.[0-9]+)?)", value)
    if not match:
        raise ValueError(f"Cannot parse numeric SVG dimension from: {value!r}")
    return float(match.group(1))


def get_svg_dimensions(svg_path: Path) -> tuple[float, float]:
    root = ET.fromstring(svg_path.read_text(encoding="utf-8"))

    width = root.attrib.get("width")
    height = root.attrib.get("height")
    if width and height:
        return parse_numeric(width), parse_numeric(height)

    view_box = root.attrib.get("viewBox")
    if view_box:
        parts = [float(part) for part in view_box.replace(",", " ").split()]
        if len(parts) == 4:
            return parts[2], parts[3]

    raise ValueError(f"Cannot determine SVG size for {svg_path}")


def px_to_emu(px: float) -> int:
    return int(px / PX_PER_INCH * EMU_PER_INCH)


def run_inkscape(svg_path: Path, emf_path: Path) -> None:
    command = [
        "inkscape",
        str(svg_path),
        "--export-type=emf",
        f"--export-filename={emf_path}",
    ]
    subprocess.run(command, check=True, capture_output=True, text=True)


def build_presentation(svg_files: list[Path], emf_dir: Path, output_path: Path) -> None:
    slide_width_px, slide_height_px = get_svg_dimensions(svg_files[0])

    presentation = Presentation()
    presentation.slide_width = Emu(px_to_emu(slide_width_px))
    presentation.slide_height = Emu(px_to_emu(slide_height_px))
    blank_layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[0]

    first_slide = presentation.slides.add_slide(blank_layout)

    for index, svg_path in enumerate(svg_files):
        slide = first_slide if index == 0 else presentation.slides.add_slide(blank_layout)
        emf_path = emf_dir / f"{svg_path.stem}.emf"
        run_inkscape(svg_path, emf_path)
        slide.shapes.add_picture(
            str(emf_path),
            left=0,
            top=0,
            # width=presentation.slide_width,
            # height=presentation.slide_height,
        )
        print(f"Added {svg_path.name} -> slide {index + 1}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))


def main() -> int:
    args = parse_args()
    input_dir = args.input_dir.resolve()
    output_path = args.output.resolve()

    Image.MAX_IMAGE_PIXELS = None
    warnings.simplefilter("ignore", Image.DecompressionBombWarning)

    if not input_dir.exists():
        print(f"Input directory does not exist: {input_dir}", file=sys.stderr)
        return 1

    svg_files = list_svg_files(input_dir)
    if not svg_files:
        print(f"No SVG files found in: {input_dir}", file=sys.stderr)
        return 1

    if args.keep_emf_dir is not None:
        emf_dir = args.keep_emf_dir.resolve()
        emf_dir.mkdir(parents=True, exist_ok=True)
        build_presentation(svg_files, emf_dir, output_path)
        print(f"Saved PPTX to {output_path}")
        print(f"Kept EMF files in {emf_dir}")
        return 0

    with TemporaryDirectory(prefix="svg2pptx_") as temp_dir:
        emf_dir = Path(temp_dir)
        build_presentation(svg_files, emf_dir, output_path)

    print(f"Saved PPTX to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
